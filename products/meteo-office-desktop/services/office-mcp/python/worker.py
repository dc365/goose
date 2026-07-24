#!/usr/bin/env python3

from __future__ import annotations

import copy
import hashlib
import io
import json
import os
import re
import shutil
import stat
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path, PurePosixPath
from typing import Any
from xml.sax.saxutils import escape

from defusedxml import ElementTree as SafeElementTree
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_BREAK
from docx.shared import Inches, Pt
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter, range_boundaries
from openpyxl.worksheet.table import Table as ExcelTable
from openpyxl.worksheet.table import TableStyleInfo
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as PresentationInches
from pptx.util import Pt as PresentationPt
from pypdf import PdfReader, PdfWriter
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.pagesizes import A4, LETTER, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import (
    Image as ReportLabImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

SCHEMA_VERSION = "meteomate.office/v1"
RUNTIME_VERSION = os.environ.get("METEOMATE_OFFICE_RUNTIME_VERSION", "1.1.0")
MAX_INPUT_BYTES = 100 * 1024 * 1024
MAX_OOXML_EXPANDED_BYTES = 500 * 1024 * 1024
MAX_OOXML_ENTRIES = 10_000
MAX_OOXML_RATIO = 200
MAX_PDF_PAGES = 1_000
MAX_RENDER_PAGES = 300
MAX_PRESENTATION_SLIDES = 300
MAX_PRESENTATION_SHAPES = 5_000
MAX_WORKSHEETS = 200
MAX_WORKSHEET_CELLS = 2_000_000
MAX_IMAGE_PIXELS = 80_000_000
ALLOWED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf"}
OOXML_MACRO_MARKERS = (
    "vbaproject.bin",
    "activex/",
    "externallinks/",
)
RELATIONSHIP_EXTERNAL = re.compile(
    br"\bTargetMode\s*=\s*['\"]External['\"]",
    re.IGNORECASE,
)
XML_DANGEROUS = re.compile(br"<!\s*(?:DOCTYPE|ENTITY)\b", re.IGNORECASE)
SPREADSHEET_FORMULA_DANGEROUS = re.compile(
    r"(?:https?|file):|\\\\|(?:WEBSERVICE|HYPERLINK|RTD|DDE|FILTERXML|ENCODEURL)\s*\(|\[[^\]]+\][^!]*!",
    re.IGNORECASE,
)


class OfficeError(Exception):
    def __init__(self, code: str, message: str):
        super().__init__(f"{code}: {message}")
        self.code = code


def fail(code: str, message: str) -> None:
    raise OfficeError(code, message)


def workspace_root() -> Path:
    raw = os.environ.get("METEOMATE_OFFICE_WORKSPACE", "")
    if not raw:
        fail("RUNTIME_UNAVAILABLE", "Office Worker 未绑定项目工作区")
    root = Path(raw)
    if not root.is_absolute() or not root.is_dir():
        fail("RUNTIME_UNAVAILABLE", "Office Worker 的项目工作区不可用")
    return root.resolve(strict=True)


WORKSPACE = workspace_root()


def inside_workspace(candidate: Path) -> bool:
    try:
        return os.path.commonpath([str(WORKSPACE), str(candidate)]) == str(WORKSPACE)
    except ValueError:
        return False


def normalize_relative_path(value: Any, label: str) -> Path:
    text = str(value or "")
    if not text or "\x00" in text:
        fail("INVALID_ARGUMENT", f"{label} 不能为空")
    if Path(text).is_absolute() or re.match(r"^[A-Za-z]:[\\/]", text):
        fail("WORKSPACE_VIOLATION", f"{label} 必须是工作区相对路径")
    parts = PurePosixPath(text.replace("\\", "/")).parts
    if ".." in parts:
        fail("WORKSPACE_VIOLATION", f"{label} 不能包含 ..")
    return Path(*parts)


def resolve_source(value: Any, label: str = "sourcePath", extensions: set[str] | None = None) -> Path:
    relative = normalize_relative_path(value, label)
    candidate = (WORKSPACE / relative).resolve(strict=True)
    if not inside_workspace(candidate):
        fail("WORKSPACE_VIOLATION", f"{label} 已超出项目工作区")
    if not candidate.is_file():
        fail("INVALID_ARGUMENT", f"{label} 不是文件")
    if extensions and candidate.suffix.lower() not in extensions:
        fail("INVALID_ARGUMENT", f"{label} 的文件格式不受支持")
    size = candidate.stat().st_size
    if size > MAX_INPUT_BYTES:
        fail("RESOURCE_LIMIT", f"{label} 超过 {MAX_INPUT_BYTES // (1024 * 1024)} MiB")
    return candidate


def resolve_output(value: Any, extension: str) -> Path:
    relative = normalize_relative_path(value, "outputPath")
    if relative.suffix.lower() != extension:
        fail("INVALID_ARGUMENT", f"outputPath 必须使用 {extension} 扩展名")
    candidate = WORKSPACE / relative
    parent = candidate.parent.resolve(strict=False)
    if not inside_workspace(parent):
        fail("WORKSPACE_VIOLATION", "outputPath 已超出项目工作区")
    parent.mkdir(parents=True, exist_ok=True, mode=0o700)
    os.chmod(parent, stat.S_IRWXU)
    resolved_parent = parent.resolve(strict=True)
    if not inside_workspace(resolved_parent):
        fail("WORKSPACE_VIOLATION", "outputPath 父目录通过符号链接逃逸")
    target = resolved_parent / candidate.name
    if target.exists():
        fail("OUTPUT_EXISTS", "输出文件已存在；请使用新的版本文件名")
    return target


def relative_path(target: Path) -> str:
    return target.resolve(strict=False).relative_to(WORKSPACE).as_posix()


def sha256_file(target: Path) -> str:
    digest = hashlib.sha256()
    with target.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def expected_hash(value: Any) -> str:
    normalized = str(value or "").lower().removeprefix("sha256:")
    if not re.fullmatch(r"[a-f0-9]{64}", normalized):
        fail("INVALID_ARGUMENT", "sourceHash 必须是 SHA-256")
    return normalized


def verify_source_hash(target: Path, value: Any) -> str:
    actual = sha256_file(target)
    if actual != expected_hash(value):
        fail("SOURCE_CHANGED", "源文件内容已变化，请重新检查后再编辑")
    return actual


def atomic_save(save: Any, target: Path) -> None:
    temporary = target.with_name(f".{target.name}.tmp-{os.getpid()}")
    try:
        save(temporary)
        os.chmod(temporary, stat.S_IRUSR | stat.S_IWUSR)
        temporary.replace(target)
    finally:
        temporary.unlink(missing_ok=True)


def media_type(target: Path) -> str:
    return {
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pdf": "application/pdf",
        ".png": "image/png",
        ".json": "application/json",
    }.get(target.suffix.lower(), "application/octet-stream")


def artifact_type(target: Path) -> str:
    return {
        ".docx": "DOCUMENT",
        ".pptx": "PRESENTATION",
        ".xlsx": "SPREADSHEET",
        ".pdf": "PDF",
        ".png": "IMAGE",
    }.get(target.suffix.lower(), "FILE")


def artifact_record(
    target: Path,
    *,
    status: str = "draft",
    metadata: dict[str, Any] | None = None,
    source_hash: str | None = None,
) -> dict[str, Any]:
    digest = sha256_file(target)
    return {
        "apiVersion": "meteomate/v1",
        "kind": "Artifact",
        "id": f"artifact-office-{digest[:24]}",
        "name": target.name,
        "type": artifact_type(target),
        "path": relative_path(target),
        "mediaType": media_type(target),
        "status": status,
        "sizeBytes": target.stat().st_size,
        "contentHash": digest,
        "metadata": {
            "source": "office-artifacts",
            "format": target.suffix.lower().removeprefix("."),
            "runtimeVersion": RUNTIME_VERSION,
            **({"sourceHash": source_hash} if source_hash else {}),
            **(metadata or {}),
        },
    }


def base_result(**values: Any) -> dict[str, Any]:
    return {"schemaVersion": SCHEMA_VERSION, **values}


def validate_embedded_chart_workbook(data: bytes) -> dict[str, int]:
    stream = io.BytesIO(data)
    if not zipfile.is_zipfile(stream):
        fail("SECURITY_REJECTED", "PPTX 图表嵌入数据不是有效的 XLSX")
    expanded_bytes = 0
    names: set[str] = set()
    with zipfile.ZipFile(stream) as archive:
        entries = archive.infolist()
        if len(entries) > MAX_OOXML_ENTRIES:
            fail("RESOURCE_LIMIT", "PPTX 图表嵌入数据 ZIP 条目数超过限制")
        for entry in entries:
            normalized = PurePosixPath(entry.filename)
            canonical = normalized.as_posix().lower()
            if entry.filename.startswith(("/", "\\")) or ".." in normalized.parts:
                fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含路径逃逸条目")
            if canonical in names:
                fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含重复 ZIP 条目")
            names.add(canonical)
            if entry.flag_bits & 0x1:
                fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含加密 ZIP 条目")
            if any(marker in canonical for marker in (*OOXML_MACRO_MARKERS, "embeddings/")):
                fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含宏、外链或嵌入对象")
            expanded_bytes += entry.file_size
            if expanded_bytes > MAX_OOXML_EXPANDED_BYTES:
                fail("RESOURCE_LIMIT", "PPTX 图表嵌入数据解压后大小超过限制")
            if entry.compress_size and entry.file_size / entry.compress_size > MAX_OOXML_RATIO:
                fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含异常压缩比条目")
            if canonical.endswith((".xml", ".rels")) and entry.file_size <= 4 * 1024 * 1024:
                content = archive.read(entry)
                if XML_DANGEROUS.search(content):
                    fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含 DTD 或实体声明")
                if canonical.endswith(".rels") and RELATIONSHIP_EXTERNAL.search(content):
                    fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含外部关系")
                if canonical.startswith("xl/worksheets/"):
                    try:
                        root = SafeElementTree.fromstring(content)
                    except Exception:
                        fail("SECURITY_REJECTED", "PPTX 图表嵌入数据包含无效工作表 XML")
                    for element in root.iter():
                        if str(element.tag).endswith("}f") and element.text:
                            validate_spreadsheet_formula(f"={element.text}")
    if not {"[content_types].xml", "xl/workbook.xml"}.issubset(names):
        fail("SECURITY_REJECTED", "PPTX 图表嵌入数据缺少 XLSX 必需条目")
    return {"entryCount": len(names), "expandedBytes": expanded_bytes}


def ooxml_preflight(target: Path) -> dict[str, Any]:
    extension = target.suffix.lower()
    required_entries = {
        ".docx": {"[content_types].xml", "word/document.xml"},
        ".pptx": {"[content_types].xml", "ppt/presentation.xml"},
        ".xlsx": {"[content_types].xml", "xl/workbook.xml"},
    }
    if extension not in required_entries:
        fail("INVALID_ARGUMENT", "OOXML 检查只支持 DOCX、PPTX 和 XLSX")
    format_name = extension.removeprefix(".").upper()
    if not zipfile.is_zipfile(target):
        fail("VALIDATION_FAILED", f"{format_name} 不是有效的 OOXML ZIP")
    total = 0
    names: set[str] = set()
    external_relationships: list[str] = []
    embedded_chart_workbooks: list[dict[str, Any]] = []
    with zipfile.ZipFile(target) as archive:
        entries = archive.infolist()
        if len(entries) > MAX_OOXML_ENTRIES:
            fail("RESOURCE_LIMIT", f"{format_name} ZIP 条目数超过限制")
        for entry in entries:
            normalized = PurePosixPath(entry.filename)
            if entry.filename.startswith(("/", "\\")) or ".." in normalized.parts:
                fail("SECURITY_REJECTED", f"{format_name} 包含路径逃逸条目")
            canonical = normalized.as_posix().lower()
            if canonical in names:
                fail("SECURITY_REJECTED", f"{format_name} 包含重复 ZIP 条目")
            names.add(canonical)
            if entry.flag_bits & 0x1:
                fail("SECURITY_REJECTED", f"{format_name} 包含加密 ZIP 条目")
            if any(marker in canonical for marker in OOXML_MACRO_MARKERS):
                fail("SECURITY_REJECTED", f"{format_name} 包含宏、ActiveX、外部链接或嵌入对象")
            total += entry.file_size
            if total > MAX_OOXML_EXPANDED_BYTES:
                fail("RESOURCE_LIMIT", f"{format_name} 解压后大小超过限制")
            if entry.compress_size and entry.file_size / entry.compress_size > MAX_OOXML_RATIO:
                fail("SECURITY_REJECTED", f"{format_name} 包含异常压缩比条目")
            if "/embeddings/" in canonical:
                if (
                    extension != ".pptx"
                    or not canonical.startswith("ppt/embeddings/")
                    or not canonical.endswith(".xlsx")
                ):
                    fail("SECURITY_REJECTED", f"{format_name} 包含不受支持的嵌入对象")
                embedded_chart_workbooks.append({
                    "path": entry.filename,
                    **validate_embedded_chart_workbook(archive.read(entry)),
                })
            if canonical.endswith((".xml", ".rels")) and entry.file_size <= 4 * 1024 * 1024:
                data = archive.read(entry)
                if XML_DANGEROUS.search(data):
                    fail("SECURITY_REJECTED", f"{format_name} XML 包含 DTD 或实体声明")
                if canonical.endswith(".rels") and RELATIONSHIP_EXTERNAL.search(data):
                    external_relationships.append(entry.filename)
    if not required_entries[extension].issubset(names):
        fail("VALIDATION_FAILED", f"{format_name} 缺少必需的 OOXML 条目")
    if external_relationships:
        fail("SECURITY_REJECTED", f"{format_name} 包含外部关系")
    return {
        "entryCount": len(names),
        "expandedBytes": total,
        "externalRelationships": [],
        "embeddedChartWorkbooks": embedded_chart_workbooks,
    }


def document_anchor_names(target: Path) -> list[dict[str, str]]:
    anchors: list[dict[str, str]] = []
    with zipfile.ZipFile(target) as archive:
        for name in archive.namelist():
            if not name.startswith("word/") or not name.endswith(".xml"):
                continue
            data = archive.read(name).decode("utf-8", errors="ignore")
            for bookmark in re.findall(r"<w:bookmarkStart\b[^>]*\bw:name=\"([^\"]+)\"", data):
                if not bookmark.startswith("_"):
                    anchors.append({"id": bookmark, "kind": "bookmark", "part": name})
            for tag in re.findall(r"<w:tag\b[^>]*\bw:val=\"([^\"]+)\"", data):
                anchors.append({"id": tag, "kind": "content-control", "part": name})
    unique = {(item["kind"], item["id"], item["part"]): item for item in anchors}
    return list(unique.values())


def iter_document_paragraphs(document: Document):
    for paragraph in document.paragraphs:
        yield paragraph
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs
    for section in document.sections:
        yield from section.header.paragraphs
        yield from section.footer.paragraphs


def replace_paragraph_text(paragraph: Any, old: str, new: str) -> int:
    if old not in paragraph.text:
        return 0
    count = paragraph.text.count(old)
    if len(paragraph.runs) == 1:
        paragraph.runs[0].text = paragraph.runs[0].text.replace(old, new)
        return count
    original = paragraph.text
    paragraph.clear()
    paragraph.add_run(original.replace(old, new))
    return count


def replace_text(document: Document, old: str, new: str) -> int:
    return sum(replace_paragraph_text(paragraph, old, new) for paragraph in iter_document_paragraphs(document))


def fill_content_control(document: Document, name: str, value: str) -> bool:
    namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    changed = False
    for control in document.element.iter(f"{namespace}sdt"):
        properties = control.find(f"{namespace}sdtPr")
        if properties is None:
            continue
        tag = properties.find(f"{namespace}tag")
        alias = properties.find(f"{namespace}alias")
        identifiers = {
            element.get(f"{namespace}val")
            for element in (tag, alias)
            if element is not None and element.get(f"{namespace}val")
        }
        if name not in identifiers:
            continue
        content = control.find(f"{namespace}sdtContent")
        if content is None:
            continue
        texts = list(content.iter(f"{namespace}t"))
        if not texts:
            continue
        texts[0].text = value
        for text in texts[1:]:
            text.text = ""
        changed = True
    return changed


def fill_bookmark(document: Document, name: str, value: str) -> bool:
    namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    for bookmark in document.element.iter(f"{namespace}bookmarkStart"):
        if bookmark.get(f"{namespace}name") != name:
            continue
        parent = bookmark.getparent()
        texts = list(parent.iter(f"{namespace}t"))
        if not texts:
            return False
        texts[0].text = value
        for text in texts[1:]:
            text.text = ""
        return True
    return False


def fill_anchor(document: Document, name: str, value: Any) -> bool:
    text = str(value)
    return fill_content_control(document, name, text) or fill_bookmark(document, name, text)


def validate_image(target: Path) -> None:
    with Image.open(target) as image:
        width, height = image.size
        if width * height > MAX_IMAGE_PIXELS:
            fail("RESOURCE_LIMIT", f"图片 {target.name} 像素数超过限制")
        image.verify()


def document_add_image(document: Document, block: dict[str, Any]) -> None:
    target = resolve_source(block.get("path"), "image.path")
    validate_image(target)
    width = block.get("widthInches")
    height = block.get("heightInches")
    kwargs: dict[str, Any] = {}
    if width is not None:
        kwargs["width"] = Inches(float(width))
    if height is not None:
        kwargs["height"] = Inches(float(height))
    document.add_picture(str(target), **kwargs)


def document_add_table(document: Document, block: dict[str, Any]) -> None:
    rows = block.get("rows")
    if not isinstance(rows, list) or not rows or not all(isinstance(row, list) for row in rows):
        fail("INVALID_ARGUMENT", "table.rows 必须是非空二维数组")
    column_count = max(len(row) for row in rows)
    table = document.add_table(rows=len(rows), cols=column_count)
    style = block.get("style")
    if style:
        table.style = str(style)
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = str(value)


def document_add_blocks(document: Document, blocks: Any) -> None:
    if not isinstance(blocks, list):
        fail("INVALID_ARGUMENT", "spec.blocks 必须是数组")
    for block in blocks:
        if not isinstance(block, dict):
            fail("INVALID_ARGUMENT", "文档内容块必须是对象")
        kind = str(block.get("type") or "paragraph")
        if kind == "paragraph":
            paragraph = document.add_paragraph(str(block.get("text") or ""))
            if block.get("style"):
                paragraph.style = str(block["style"])
        elif kind == "heading":
            document.add_heading(str(block.get("text") or ""), level=int(block.get("level") or 1))
        elif kind == "table":
            document_add_table(document, block)
        elif kind == "image":
            document_add_image(document, block)
        elif kind == "page_break":
            document.add_page_break()
        elif kind == "spacer":
            document.add_paragraph("")
        else:
            fail("UNSUPPORTED_FEATURE", f"DOCX 不支持内容块 {kind}")


def configure_document(document: Document, spec: dict[str, Any]) -> None:
    font_name = str(spec.get("defaultFont") or "")
    font_size = spec.get("defaultFontSize")
    if font_name or font_size:
        normal = document.styles["Normal"]
        if font_name:
            normal.font.name = font_name
        if font_size:
            normal.font.size = Pt(float(font_size))
    page = spec.get("page")
    if isinstance(page, dict):
        for section in document.sections:
            if str(page.get("orientation") or "").lower() == "landscape":
                section.orientation = WD_ORIENT.LANDSCAPE
                section.page_width, section.page_height = section.page_height, section.page_width
            for key, attribute in (
                ("topMarginInches", "top_margin"),
                ("bottomMarginInches", "bottom_margin"),
                ("leftMarginInches", "left_margin"),
                ("rightMarginInches", "right_margin"),
            ):
                if page.get(key) is not None:
                    setattr(section, attribute, Inches(float(page[key])))
    if spec.get("header") is not None:
        for section in document.sections:
            section.header.paragraphs[0].text = str(spec["header"])
    if spec.get("footer") is not None:
        for section in document.sections:
            section.footer.paragraphs[0].text = str(spec["footer"])


def docx_inspect(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".docx"})
    security = ooxml_preflight(source)
    document = Document(str(source))
    paragraphs = [
        {
            "id": f"paragraph:{index}",
            "text": paragraph.text[:4_000],
            "style": paragraph.style.name if paragraph.style else None,
        }
        for index, paragraph in enumerate(document.paragraphs[:2_000])
    ]
    tables = [
        {
            "id": f"table:{index}",
            "rows": len(table.rows),
            "columns": len(table.columns),
            "preview": [[cell.text[:200] for cell in row.cells[:20]] for row in table.rows[:20]],
        }
        for index, table in enumerate(document.tables[:200])
    ]
    fonts = sorted({
        run.font.name
        for paragraph in iter_document_paragraphs(document)
        for run in paragraph.runs
        if run.font.name
    })
    with zipfile.ZipFile(source) as archive:
        media = [
            {"path": name, "sizeBytes": archive.getinfo(name).file_size}
            for name in archive.namelist()
            if name.startswith("word/media/")
        ][:500]
    return base_result(
        format="docx",
        sourcePath=relative_path(source),
        sourceHash=sha256_file(source),
        structure={
            "sections": len(document.sections),
            "paragraphs": paragraphs,
            "tables": tables,
            "styles": [style.name for style in list(document.styles)[:500]],
        },
        anchors=document_anchor_names(source),
        fonts=fonts,
        media=media,
        security=security,
        warnings=[],
    )


def docx_create(payload: dict[str, Any]) -> dict[str, Any]:
    output = resolve_output(payload.get("outputPath"), ".docx")
    spec = payload.get("spec")
    if not isinstance(spec, dict):
        fail("INVALID_ARGUMENT", "spec 必须是对象")
    template_path = payload.get("templatePath")
    template_id = None
    template_hash = None
    if template_path:
        template = resolve_source(template_path, "templatePath", {".docx"})
        ooxml_preflight(template)
        template_hash = sha256_file(template)
        if payload.get("templateHash") and template_hash != expected_hash(payload["templateHash"]):
            fail("SOURCE_CHANGED", "模板内容与 templateHash 不一致")
        document = Document(str(template))
        template_id = Path(str(template_path)).stem
    else:
        document = Document()
    configure_document(document, spec)
    anchors = spec.get("anchors")
    if anchors is not None:
        if not isinstance(anchors, dict):
            fail("INVALID_ARGUMENT", "spec.anchors 必须是对象")
        missing = [name for name, value in anchors.items() if not fill_anchor(document, str(name), value)]
        if missing:
            fail("VALIDATION_FAILED", f"模板缺少锚点：{', '.join(missing[:20])}")
    if spec.get("title") and not template_path:
        document.add_heading(str(spec["title"]), level=0)
    if "blocks" in spec:
        document_add_blocks(document, spec["blocks"])
    atomic_save(document.save, output)
    ooxml_preflight(output)
    return base_result(
        artifact=artifact_record(
            output,
            metadata={
                "templateId": template_id,
                "templateHash": template_hash,
            },
        ),
        warnings=[],
    )


def docx_edit(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".docx"})
    source_hash = verify_source_hash(source, payload.get("sourceHash"))
    ooxml_preflight(source)
    output = resolve_output(payload.get("outputPath"), ".docx")
    operations = payload.get("operations")
    if not isinstance(operations, list) or not operations:
        fail("INVALID_ARGUMENT", "operations 必须是非空数组")
    document = Document(str(source))
    operation_results: list[dict[str, Any]] = []
    for index, operation in enumerate(operations):
        if not isinstance(operation, dict):
            fail("INVALID_ARGUMENT", f"operations[{index}] 必须是对象")
        name = str(operation.get("op") or "")
        if name == "replace_text":
            old = str(operation.get("old") or "")
            if not old:
                fail("INVALID_ARGUMENT", "replace_text.old 不能为空")
            count = replace_text(document, old, str(operation.get("new") or ""))
            if not count and operation.get("required", True):
                fail("VALIDATION_FAILED", f"未找到待替换文本：{old[:80]}")
            operation_results.append({"index": index, "op": name, "replacements": count})
        elif name == "set_anchor":
            anchor = str(operation.get("anchor") or "")
            if not anchor or not fill_anchor(document, anchor, operation.get("value", "")):
                fail("VALIDATION_FAILED", f"未找到模板锚点：{anchor}")
            operation_results.append({"index": index, "op": name, "anchor": anchor})
        elif name == "append_paragraph":
            paragraph = document.add_paragraph(str(operation.get("text") or ""))
            if operation.get("style"):
                paragraph.style = str(operation["style"])
            operation_results.append({"index": index, "op": name})
        elif name == "add_heading":
            document.add_heading(str(operation.get("text") or ""), level=int(operation.get("level") or 1))
            operation_results.append({"index": index, "op": name})
        elif name == "add_table":
            document_add_table(document, operation)
            operation_results.append({"index": index, "op": name})
        elif name == "add_image":
            document_add_image(document, operation)
            operation_results.append({"index": index, "op": name})
        elif name == "page_break":
            document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
            operation_results.append({"index": index, "op": name})
        elif name in {"set_header", "set_footer"}:
            attribute = "header" if name == "set_header" else "footer"
            for section in document.sections:
                getattr(section, attribute).paragraphs[0].text = str(operation.get("text") or "")
            operation_results.append({"index": index, "op": name})
        else:
            fail("UNSUPPORTED_FEATURE", f"DOCX 不支持操作 {name}")
    atomic_save(document.save, output)
    ooxml_preflight(output)
    return base_result(
        artifact=artifact_record(
            output,
            metadata={"sourcePath": relative_path(source)},
            source_hash=source_hash,
        ),
        operations=operation_results,
        warnings=[],
    )


def presentation_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text
    if getattr(shape, "has_table", False):
        return "\n".join(
            cell.text
            for row in shape.table.rows
            for cell in row.cells
            if cell.text
        )
    return ""


def presentation_shapes(presentation: Presentation):
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            yield slide_number, slide, shape


def presentation_shape_record(slide_number: int, shape: Any) -> dict[str, Any]:
    record: dict[str, Any] = {
        "id": f"slide:{slide_number}/shape:{shape.shape_id}",
        "name": shape.name,
        "type": getattr(shape.shape_type, "name", str(shape.shape_type)),
        "leftInches": round(shape.left / 914400, 4),
        "topInches": round(shape.top / 914400, 4),
        "widthInches": round(shape.width / 914400, 4),
        "heightInches": round(shape.height / 914400, 4),
    }
    text = presentation_text(shape)
    if text:
        record["textPreview"] = text[:4_000]
    if getattr(shape, "is_placeholder", False):
        record["placeholder"] = {
            "index": shape.placeholder_format.idx,
            "type": getattr(shape.placeholder_format.type, "name", str(shape.placeholder_format.type)),
        }
    if getattr(shape, "has_table", False):
        record["table"] = {
            "rows": len(shape.table.rows),
            "columns": len(shape.table.columns),
            "preview": [
                [cell.text[:200] for cell in list(row.cells)[:20]]
                for row in list(shape.table.rows)[:20]
            ],
        }
    if getattr(shape, "has_chart", False):
        record["chart"] = {
            "type": getattr(shape.chart.chart_type, "name", str(shape.chart.chart_type)),
            "series": [
                {
                    "name": str(series.name or ""),
                    "valueCount": len(series.values),
                }
                for series in shape.chart.series
            ],
        }
    return record


def find_presentation_shape(
    presentation: Presentation,
    name: Any,
    slide_number: Any = None,
) -> tuple[Any, Any]:
    target_name = str(name or "")
    if not target_name:
        fail("INVALID_ARGUMENT", "shapeName 不能为空")
    requested_slide = int(slide_number) if slide_number is not None else None
    matches = [
        (slide, shape)
        for number, slide, shape in presentation_shapes(presentation)
        if shape.name == target_name and (requested_slide is None or number == requested_slide)
    ]
    if not matches:
        fail("VALIDATION_FAILED", f"未找到命名形状：{target_name}")
    if len(matches) > 1:
        fail("INVALID_ARGUMENT", f"命名形状不唯一，请同时提供 slide：{target_name}")
    return matches[0]


def replace_presentation_text(presentation: Presentation, old: str, new: str) -> int:
    replacements = 0
    for _, _, shape in presentation_shapes(presentation):
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    count = run.text.count(old)
                    if count:
                        run.text = run.text.replace(old, new)
                        replacements += count
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    count = cell.text.count(old)
                    if count:
                        cell.text = cell.text.replace(old, new)
                        replacements += count
    return replacements


def presentation_chart_data(value: Any) -> ChartData:
    if not isinstance(value, dict):
        fail("INVALID_ARGUMENT", "chart.data 必须是对象")
    categories = value.get("categories")
    series = value.get("series")
    if not isinstance(categories, list) or not categories:
        fail("INVALID_ARGUMENT", "chart.data.categories 必须是非空数组")
    if not isinstance(series, list) or not series:
        fail("INVALID_ARGUMENT", "chart.data.series 必须是非空数组")
    data = ChartData()
    data.categories = [str(category) for category in categories]
    for index, item in enumerate(series):
        if not isinstance(item, dict):
            fail("INVALID_ARGUMENT", f"chart.data.series[{index}] 必须是对象")
        values = item.get("values")
        if not isinstance(values, list) or len(values) != len(categories):
            fail("INVALID_ARGUMENT", f"chart.data.series[{index}].values 长度必须与 categories 一致")
        try:
            numeric_values = [None if value is None else float(value) for value in values]
        except (TypeError, ValueError):
            fail("INVALID_ARGUMENT", f"chart.data.series[{index}].values 必须是数值数组")
        data.add_series(str(item.get("name") or f"系列 {index + 1}"), numeric_values)
    return data


def presentation_chart_type(value: Any) -> Any:
    chart_types = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    normalized = str(value or "column").lower()
    if normalized not in chart_types:
        fail("UNSUPPORTED_FEATURE", f"PPTX 不支持图表类型 {normalized}")
    return chart_types[normalized]


def presentation_inches(value: Any, default: float) -> Any:
    try:
        return PresentationInches(float(default if value is None else value))
    except (TypeError, ValueError):
        fail("INVALID_ARGUMENT", "PPTX 元素位置和尺寸必须是数字")


def set_presentation_text(shape: Any, value: Any) -> None:
    if not getattr(shape, "has_text_frame", False):
        fail("UNSUPPORTED_FEATURE", f"形状 {shape.name} 不支持文本")
    shape.text = str(value or "")


def set_presentation_table(shape: Any, rows: Any) -> None:
    if not getattr(shape, "has_table", False):
        fail("UNSUPPORTED_FEATURE", f"形状 {shape.name} 不是表格")
    if not isinstance(rows, list) or not rows or not all(isinstance(row, list) for row in rows):
        fail("INVALID_ARGUMENT", "table.rows 必须是非空二维数组")
    table = shape.table
    if len(rows) > len(table.rows) or max(len(row) for row in rows) > len(table.columns):
        fail("UNSUPPORTED_FEATURE", "PPTX 表格更新不能超过模板现有行列")
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = str(value)


def replace_presentation_image(slide: Any, shape: Any, image_path: Any) -> Any:
    target = resolve_source(image_path, "image.path")
    validate_image(target)
    left, top, width, height, name = shape.left, shape.top, shape.width, shape.height, shape.name
    element = shape._element
    element.getparent().remove(element)
    replacement = slide.shapes.add_picture(str(target), left, top, width=width, height=height)
    replacement.name = name
    return replacement


def presentation_layout(presentation: Presentation, value: Any) -> Any:
    if value is not None:
        if isinstance(value, int) or str(value).isdigit():
            index = int(value)
            if index < 0 or index >= len(presentation.slide_layouts):
                fail("INVALID_ARGUMENT", f"PPTX layout 索引超出范围：{index}")
            return presentation.slide_layouts[index]
        for layout in presentation.slide_layouts:
            if layout.name == str(value):
                return layout
        fail("VALIDATION_FAILED", f"PPTX 模板缺少布局：{value}")
    blank = next((layout for layout in presentation.slide_layouts if layout.name.lower() == "blank"), None)
    return blank or presentation.slide_layouts[min(6, len(presentation.slide_layouts) - 1)]


def presentation_add_text(slide: Any, element: dict[str, Any]) -> Any:
    shape = slide.shapes.add_textbox(
        presentation_inches(element.get("x"), 0.7),
        presentation_inches(element.get("y"), 0.7),
        presentation_inches(element.get("width"), 8.6),
        presentation_inches(element.get("height"), 1.0),
    )
    shape.text = str(element.get("text") or "")
    if element.get("name"):
        shape.name = str(element["name"])
    for paragraph in shape.text_frame.paragraphs:
        paragraph.level = max(0, min(8, int(element.get("level") or 0)))
        paragraph.alignment = None
        for run in paragraph.runs:
            if element.get("fontSize") is not None:
                run.font.size = PresentationPt(float(element["fontSize"]))
            if element.get("font"):
                run.font.name = str(element["font"])
            if element.get("bold") is not None:
                run.font.bold = bool(element["bold"])
            color = str(element.get("color") or "").removeprefix("#")
            if re.fullmatch(r"[a-fA-F0-9]{6}", color):
                run.font.color.rgb = RGBColor.from_string(color.upper())
    return shape


def presentation_add_table(slide: Any, element: dict[str, Any]) -> Any:
    rows = element.get("rows")
    if not isinstance(rows, list) or not rows or not all(isinstance(row, list) for row in rows):
        fail("INVALID_ARGUMENT", "table.rows 必须是非空二维数组")
    column_count = max(len(row) for row in rows)
    shape = slide.shapes.add_table(
        len(rows),
        column_count,
        presentation_inches(element.get("x"), 0.7),
        presentation_inches(element.get("y"), 1.8),
        presentation_inches(element.get("width"), 8.6),
        presentation_inches(element.get("height"), 4.8),
    )
    if element.get("name"):
        shape.name = str(element["name"])
    set_presentation_table(shape, rows)
    return shape


def presentation_add_chart(slide: Any, element: dict[str, Any]) -> Any:
    shape = slide.shapes.add_chart(
        presentation_chart_type(element.get("chartType")),
        presentation_inches(element.get("x"), 0.7),
        presentation_inches(element.get("y"), 1.8),
        presentation_inches(element.get("width"), 8.6),
        presentation_inches(element.get("height"), 4.8),
        presentation_chart_data(element.get("data")),
    )
    if element.get("name"):
        shape.name = str(element["name"])
    chart = shape.chart
    if element.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = str(element["title"])
    chart.has_legend = bool(element.get("legend", True))
    return shape


def presentation_add_slide(presentation: Presentation, spec: Any) -> Any:
    if not isinstance(spec, dict):
        fail("INVALID_ARGUMENT", "slide 必须是对象")
    slide = presentation.slides.add_slide(presentation_layout(presentation, spec.get("layout")))
    title = spec.get("title")
    if title is not None:
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(title)
        else:
            presentation_add_text(slide, {
                "name": "Title",
                "text": title,
                "x": 0.7,
                "y": 0.35,
                "width": 8.6,
                "height": 0.8,
                "fontSize": 26,
                "bold": True,
            })
    subtitle = spec.get("subtitle")
    if subtitle is not None:
        placeholder = next(
            (
                shape
                for shape in slide.placeholders
                if getattr(shape.placeholder_format.type, "name", "") == "SUBTITLE"
            ),
            None,
        )
        if placeholder is not None:
            placeholder.text = str(subtitle)
        else:
            presentation_add_text(slide, {
                "name": "Subtitle",
                "text": subtitle,
                "x": 0.7,
                "y": 1.15,
                "width": 8.6,
                "height": 0.6,
                "fontSize": 15,
            })
    elements = spec.get("elements", [])
    if not isinstance(elements, list):
        fail("INVALID_ARGUMENT", "slide.elements 必须是数组")
    for index, element in enumerate(elements):
        if not isinstance(element, dict):
            fail("INVALID_ARGUMENT", f"slide.elements[{index}] 必须是对象")
        kind = str(element.get("type") or "text")
        if kind == "text":
            presentation_add_text(slide, element)
        elif kind == "image":
            target = resolve_source(element.get("path"), "image.path")
            validate_image(target)
            shape = slide.shapes.add_picture(
                str(target),
                presentation_inches(element.get("x"), 0.7),
                presentation_inches(element.get("y"), 1.8),
                width=presentation_inches(element.get("width"), 8.6),
                height=presentation_inches(element.get("height"), 4.8),
            )
            if element.get("name"):
                shape.name = str(element["name"])
        elif kind == "table":
            presentation_add_table(slide, element)
        elif kind == "chart":
            presentation_add_chart(slide, element)
        else:
            fail("UNSUPPORTED_FEATURE", f"PPTX 不支持元素 {kind}")
    if spec.get("notes") is not None:
        slide.notes_slide.notes_text_frame.text = str(spec["notes"])
    return slide


def apply_presentation_anchor(presentation: Presentation, name: str, value: Any) -> None:
    slide, shape = find_presentation_shape(presentation, name)
    if not isinstance(value, dict):
        set_presentation_text(shape, value)
        return
    kind = str(value.get("type") or "text")
    if kind == "text":
        set_presentation_text(shape, value.get("text"))
    elif kind == "image":
        replace_presentation_image(slide, shape, value.get("path"))
    elif kind == "table":
        set_presentation_table(shape, value.get("rows"))
    elif kind == "chart":
        if not getattr(shape, "has_chart", False):
            fail("UNSUPPORTED_FEATURE", f"形状 {name} 不是图表")
        shape.chart.replace_data(presentation_chart_data(value.get("data")))
    else:
        fail("UNSUPPORTED_FEATURE", f"PPTX 锚点不支持类型 {kind}")


def pptx_inspect(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".pptx"})
    security = ooxml_preflight(source)
    presentation = Presentation(str(source))
    if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
        fail("RESOURCE_LIMIT", "PPTX 幻灯片数量超过限制")
    shape_count = sum(len(slide.shapes) for slide in presentation.slides)
    if shape_count > MAX_PRESENTATION_SHAPES:
        fail("RESOURCE_LIMIT", "PPTX 形状数量超过限制")
    slides = [
        {
            "id": f"slide:{slide_number}",
            "number": slide_number,
            "layout": slide.slide_layout.name,
            "shapes": [
                presentation_shape_record(slide_number, shape)
                for shape in slide.shapes
            ],
            "notesPreview": slide.notes_slide.notes_text_frame.text[:4_000]
            if slide.has_notes_slide
            else "",
        }
        for slide_number, slide in enumerate(presentation.slides, start=1)
    ]
    with zipfile.ZipFile(source) as archive:
        media = [
            {"path": name, "sizeBytes": archive.getinfo(name).file_size}
            for name in archive.namelist()
            if name.startswith("ppt/media/")
        ][:500]
    return base_result(
        format="pptx",
        sourcePath=relative_path(source),
        sourceHash=sha256_file(source),
        slideCount=len(slides),
        slideSize={
            "widthInches": round(presentation.slide_width / 914400, 4),
            "heightInches": round(presentation.slide_height / 914400, 4),
        },
        slides=slides,
        anchors=[
            {
                "id": shape.name,
                "kind": "named-shape",
                "slide": slide_number,
                "shapeId": shape.shape_id,
            }
            for slide_number, _, shape in presentation_shapes(presentation)
        ],
        media=media,
        security=security,
        warnings=[],
    )


def pptx_create(payload: dict[str, Any]) -> dict[str, Any]:
    output = resolve_output(payload.get("outputPath"), ".pptx")
    spec = payload.get("spec")
    if not isinstance(spec, dict):
        fail("INVALID_ARGUMENT", "spec 必须是对象")
    template_path = payload.get("templatePath")
    template_id = None
    template_hash = None
    if template_path:
        template = resolve_source(template_path, "templatePath", {".pptx"})
        ooxml_preflight(template)
        template_hash = sha256_file(template)
        if payload.get("templateHash") and template_hash != expected_hash(payload["templateHash"]):
            fail("SOURCE_CHANGED", "模板内容与 templateHash 不一致")
        presentation = Presentation(str(template))
        template_id = Path(str(template_path)).stem
    else:
        presentation = Presentation()
    if spec.get("layout") == "wide":
        presentation.slide_width = PresentationInches(13.333333)
        presentation.slide_height = PresentationInches(7.5)
    anchors = spec.get("anchors")
    if anchors is not None:
        if not isinstance(anchors, dict):
            fail("INVALID_ARGUMENT", "spec.anchors 必须是对象")
        for name, value in anchors.items():
            apply_presentation_anchor(presentation, str(name), value)
    slides = spec.get("slides", [])
    if not isinstance(slides, list):
        fail("INVALID_ARGUMENT", "spec.slides 必须是数组")
    for slide in slides:
        presentation_add_slide(presentation, slide)
    if not presentation.slides:
        fail("VALIDATION_FAILED", "PPTX 至少需要一页幻灯片")
    if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
        fail("RESOURCE_LIMIT", "PPTX 幻灯片数量超过限制")
    atomic_save(presentation.save, output)
    ooxml_preflight(output)
    return base_result(
        artifact=artifact_record(
            output,
            metadata={
                "templateId": template_id,
                "templateHash": template_hash,
                "slideCount": len(presentation.slides),
            },
        ),
        warnings=[],
    )


def pptx_edit(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".pptx"})
    source_hash = verify_source_hash(source, payload.get("sourceHash"))
    ooxml_preflight(source)
    output = resolve_output(payload.get("outputPath"), ".pptx")
    operations = payload.get("operations")
    if not isinstance(operations, list) or not operations:
        fail("INVALID_ARGUMENT", "operations 必须是非空数组")
    presentation = Presentation(str(source))
    results: list[dict[str, Any]] = []
    for index, operation in enumerate(operations):
        if not isinstance(operation, dict):
            fail("INVALID_ARGUMENT", f"operations[{index}] 必须是对象")
        name = str(operation.get("op") or "")
        if name == "replace_text":
            old = str(operation.get("old") or "")
            if not old:
                fail("INVALID_ARGUMENT", "replace_text.old 不能为空")
            count = replace_presentation_text(presentation, old, str(operation.get("new") or ""))
            if not count and operation.get("required", True):
                fail("VALIDATION_FAILED", f"未找到待替换文本：{old[:80]}")
            results.append({"index": index, "op": name, "replacements": count})
        elif name == "set_shape_text":
            _, shape = find_presentation_shape(
                presentation,
                operation.get("shapeName"),
                operation.get("slide"),
            )
            set_presentation_text(shape, operation.get("text"))
            results.append({"index": index, "op": name, "shapeName": shape.name})
        elif name == "replace_image":
            slide, shape = find_presentation_shape(
                presentation,
                operation.get("shapeName"),
                operation.get("slide"),
            )
            replace_presentation_image(slide, shape, operation.get("path"))
            results.append({"index": index, "op": name, "shapeName": shape.name})
        elif name == "set_table":
            _, shape = find_presentation_shape(
                presentation,
                operation.get("shapeName"),
                operation.get("slide"),
            )
            set_presentation_table(shape, operation.get("rows"))
            results.append({"index": index, "op": name, "shapeName": shape.name})
        elif name == "set_chart_data":
            _, shape = find_presentation_shape(
                presentation,
                operation.get("shapeName"),
                operation.get("slide"),
            )
            if not getattr(shape, "has_chart", False):
                fail("UNSUPPORTED_FEATURE", f"形状 {shape.name} 不是图表")
            shape.chart.replace_data(presentation_chart_data(operation.get("data")))
            results.append({"index": index, "op": name, "shapeName": shape.name})
        elif name == "add_slide":
            presentation_add_slide(presentation, operation.get("slide"))
            results.append({"index": index, "op": name, "slide": len(presentation.slides)})
        elif name == "set_notes":
            slide_number = int(operation.get("slide") or 0)
            if slide_number < 1 or slide_number > len(presentation.slides):
                fail("INVALID_ARGUMENT", "set_notes.slide 超出范围")
            presentation.slides[slide_number - 1].notes_slide.notes_text_frame.text = str(
                operation.get("text") or ""
            )
            results.append({"index": index, "op": name, "slide": slide_number})
        else:
            fail("UNSUPPORTED_FEATURE", f"PPTX 不支持操作 {name}")
    if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
        fail("RESOURCE_LIMIT", "PPTX 幻灯片数量超过限制")
    atomic_save(presentation.save, output)
    ooxml_preflight(output)
    return base_result(
        artifact=artifact_record(
            output,
            metadata={
                "sourcePath": relative_path(source),
                "slideCount": len(presentation.slides),
            },
            source_hash=source_hash,
        ),
        operations=results,
        warnings=[],
    )


def spreadsheet_value(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


def validate_spreadsheet_formula(value: Any) -> None:
    formula = str(value or "")
    if formula.startswith("=") and SPREADSHEET_FORMULA_DANGEROUS.search(formula):
        fail("SECURITY_REJECTED", "XLSX 公式包含外部访问、DDE 或外部工作簿引用")


def spreadsheet_sheet_name(value: Any) -> str:
    name = str(value or "")
    if not name or len(name) > 31 or re.search(r"[\[\]:*?/\\]", name):
        fail("INVALID_ARGUMENT", f"工作表名称不合法：{name}")
    return name


def spreadsheet_sheet(workbook: Workbook, value: Any) -> Any:
    name = spreadsheet_sheet_name(value)
    if name not in workbook.sheetnames:
        fail("VALIDATION_FAILED", f"未找到工作表：{name}")
    return workbook[name]


def spreadsheet_color(value: Any, default: str = "000000") -> str:
    color = str(value or default).removeprefix("#").upper()
    if not re.fullmatch(r"[A-F0-9]{6}(?:[A-F0-9]{2})?", color):
        fail("INVALID_ARGUMENT", f"颜色值不合法：{value}")
    return f"FF{color}" if len(color) == 6 else color


def spreadsheet_style(cell: Any, style: Any) -> None:
    if not isinstance(style, dict):
        fail("INVALID_ARGUMENT", "单元格样式必须是对象")
    font = style.get("font")
    if isinstance(font, dict):
        cell.font = Font(
            name=str(font.get("name")) if font.get("name") else cell.font.name,
            size=float(font["size"]) if font.get("size") is not None else cell.font.size,
            bold=bool(font.get("bold", cell.font.bold)),
            italic=bool(font.get("italic", cell.font.italic)),
            color=spreadsheet_color(font["color"]) if font.get("color") else cell.font.color,
        )
    if style.get("fill"):
        cell.fill = PatternFill("solid", fgColor=spreadsheet_color(style["fill"], "FFFFFF"))
    alignment = style.get("alignment")
    if isinstance(alignment, dict):
        cell.alignment = Alignment(
            horizontal=str(alignment.get("horizontal")) if alignment.get("horizontal") else None,
            vertical=str(alignment.get("vertical")) if alignment.get("vertical") else None,
            wrap_text=bool(alignment.get("wrapText", False)),
        )
    if style.get("numberFormat") is not None:
        cell.number_format = str(style["numberFormat"])
    border = style.get("border")
    if isinstance(border, dict):
        side = Side(
            style=str(border.get("style") or "thin"),
            color=spreadsheet_color(border.get("color"), "B8C2CC"),
        )
        cell.border = Border(left=side, right=side, top=side, bottom=side)


def spreadsheet_range_cells(worksheet: Any, cell_range: Any):
    text = str(cell_range or "")
    try:
        min_column, min_row, max_column, max_row = range_boundaries(text)
    except ValueError:
        fail("INVALID_ARGUMENT", f"单元格区域不合法：{text}")
    if min_row < 1 or min_column < 1 or max_row * max_column > MAX_WORKSHEET_CELLS:
        fail("RESOURCE_LIMIT", f"单元格区域超出限制：{text}")
    for row in worksheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_column,
        max_col=max_column,
    ):
        yield from row


def spreadsheet_write_range(worksheet: Any, start_cell: Any, values: Any) -> None:
    if not isinstance(values, list) or not values or not all(isinstance(row, list) for row in values):
        fail("INVALID_ARGUMENT", "values 必须是非空二维数组")
    try:
        start_column, start_row, _, _ = range_boundaries(str(start_cell or "A1"))
    except ValueError:
        fail("INVALID_ARGUMENT", f"起始单元格不合法：{start_cell}")
    column_count = max((len(row) for row in values), default=0)
    if (start_row + len(values) - 1) * (start_column + column_count - 1) > MAX_WORKSHEET_CELLS:
        fail("RESOURCE_LIMIT", "工作表写入区域超过限制")
    for row_offset, row in enumerate(values):
        for column_offset, value in enumerate(row):
            validate_spreadsheet_formula(value)
            worksheet.cell(
                row=start_row + row_offset,
                column=start_column + column_offset,
                value=value,
            )


def spreadsheet_set_cells(worksheet: Any, cells: Any) -> None:
    if not isinstance(cells, dict):
        fail("INVALID_ARGUMENT", "cells 必须是单元格地址到值的对象")
    for coordinate, value in cells.items():
        try:
            cell = worksheet[str(coordinate)]
        except ValueError:
            fail("INVALID_ARGUMENT", f"单元格地址不合法：{coordinate}")
        if isinstance(value, dict):
            if "formula" in value:
                formula = str(value["formula"])
                normalized_formula = formula if formula.startswith("=") else f"={formula}"
                validate_spreadsheet_formula(normalized_formula)
                cell.value = normalized_formula
            elif "value" in value:
                cell.value = value["value"]
            if value.get("style") is not None:
                spreadsheet_style(cell, value["style"])
        else:
            cell.value = value


def spreadsheet_add_table(worksheet: Any, value: Any) -> None:
    if not isinstance(value, dict):
        fail("INVALID_ARGUMENT", "table 必须是对象")
    reference = str(value.get("ref") or "")
    list(spreadsheet_range_cells(worksheet, reference))
    name = str(value.get("name") or "")
    if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_.]{0,254}", name):
        fail("INVALID_ARGUMENT", f"Excel 表名称不合法：{name}")
    if any(name == table.displayName for table in worksheet.tables.values()):
        fail("OUTPUT_EXISTS", f"Excel 表名称已存在：{name}")
    table = ExcelTable(displayName=name, ref=reference)
    table.tableStyleInfo = TableStyleInfo(
        name=str(value.get("style") or "TableStyleMedium2"),
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=bool(value.get("rowStripes", True)),
        showColumnStripes=bool(value.get("columnStripes", False)),
    )
    worksheet.add_table(table)


def spreadsheet_chart(value: Any, worksheet: Any) -> Any:
    if not isinstance(value, dict):
        fail("INVALID_ARGUMENT", "chart 必须是对象")
    chart_type = str(value.get("chartType") or "column").lower()
    if chart_type == "column":
        chart = BarChart()
        chart.type = "col"
    elif chart_type == "bar":
        chart = BarChart()
        chart.type = "bar"
    elif chart_type == "line":
        chart = LineChart()
    elif chart_type == "pie":
        chart = PieChart()
    else:
        fail("UNSUPPORTED_FEATURE", f"XLSX 不支持图表类型 {chart_type}")
    data_range = str(value.get("dataRange") or "")
    try:
        min_column, min_row, max_column, max_row = range_boundaries(data_range)
    except ValueError:
        fail("INVALID_ARGUMENT", f"chart.dataRange 不合法：{data_range}")
    titles_from_data = bool(value.get("titlesFromData", True))
    chart.add_data(
        Reference(
            worksheet,
            min_col=min_column,
            min_row=min_row,
            max_col=max_column,
            max_row=max_row,
        ),
        titles_from_data=titles_from_data,
        from_rows=bool(value.get("seriesFromRows", False)),
    )
    categories_range = value.get("categoriesRange")
    if categories_range:
        try:
            category_column, category_row, category_max_column, category_max_row = range_boundaries(
                str(categories_range)
            )
        except ValueError:
            fail("INVALID_ARGUMENT", f"chart.categoriesRange 不合法：{categories_range}")
        chart.set_categories(
            Reference(
                worksheet,
                min_col=category_column,
                min_row=category_row,
                max_col=category_max_column,
                max_row=category_max_row,
            )
        )
    chart.title = str(value.get("title") or "")
    chart.style = int(value.get("style") or 10)
    chart.height = float(value.get("height") or 7.5)
    chart.width = float(value.get("width") or 13.0)
    if value.get("xAxisTitle"):
        chart.x_axis.title = str(value["xAxisTitle"])
    if value.get("yAxisTitle") and hasattr(chart, "y_axis"):
        chart.y_axis.title = str(value["yAxisTitle"])
    return chart


def spreadsheet_add_chart(worksheet: Any, value: Any) -> None:
    chart = spreadsheet_chart(value, worksheet)
    worksheet.add_chart(chart, str(value.get("anchor") or "F2"))


def configure_spreadsheet_calculation(workbook: Workbook) -> None:
    workbook.calculation.fullCalcOnLoad = True
    workbook.calculation.forceFullCalc = True
    workbook.calculation.calcMode = "auto"


def apply_worksheet_spec(worksheet: Any, spec: Any) -> None:
    if not isinstance(spec, dict):
        fail("INVALID_ARGUMENT", "worksheet 必须是对象")
    if spec.get("data") is not None:
        spreadsheet_write_range(worksheet, spec.get("startCell") or "A1", spec["data"])
    if spec.get("cells") is not None:
        spreadsheet_set_cells(worksheet, spec["cells"])
    styles = spec.get("styles", [])
    if not isinstance(styles, list):
        fail("INVALID_ARGUMENT", "worksheet.styles 必须是数组")
    for entry in styles:
        if not isinstance(entry, dict):
            fail("INVALID_ARGUMENT", "worksheet.styles 条目必须是对象")
        for cell in spreadsheet_range_cells(worksheet, entry.get("range")):
            spreadsheet_style(cell, entry.get("style"))
    if spec.get("freezePane") is not None:
        worksheet.freeze_panes = str(spec["freezePane"]) if spec["freezePane"] else None
    if spec.get("autoFilter") is not None:
        worksheet.auto_filter.ref = str(spec["autoFilter"])
    if spec.get("printArea") is not None:
        worksheet.print_area = str(spec["printArea"])
    if spec.get("orientation") is not None:
        orientation = str(spec["orientation"]).lower()
        if orientation not in {"portrait", "landscape"}:
            fail("INVALID_ARGUMENT", "worksheet.orientation 必须是 portrait 或 landscape")
        worksheet.page_setup.orientation = orientation
    if spec.get("fitToWidth") is not None:
        worksheet.page_setup.fitToWidth = int(spec["fitToWidth"])
        worksheet.sheet_properties.pageSetUpPr.fitToPage = True
    if spec.get("fitToHeight") is not None:
        worksheet.page_setup.fitToHeight = int(spec["fitToHeight"])
        worksheet.sheet_properties.pageSetUpPr.fitToPage = True
    widths = spec.get("columnWidths", {})
    if not isinstance(widths, dict):
        fail("INVALID_ARGUMENT", "worksheet.columnWidths 必须是对象")
    for column, width in widths.items():
        key = get_column_letter(int(column)) if str(column).isdigit() else str(column).upper()
        worksheet.column_dimensions[key].width = float(width)
    heights = spec.get("rowHeights", {})
    if not isinstance(heights, dict):
        fail("INVALID_ARGUMENT", "worksheet.rowHeights 必须是对象")
    for row, height in heights.items():
        worksheet.row_dimensions[int(row)].height = float(height)
    merged_cells = spec.get("mergedCells", [])
    if not isinstance(merged_cells, list):
        fail("INVALID_ARGUMENT", "worksheet.mergedCells 必须是数组")
    for cell_range in merged_cells:
        list(spreadsheet_range_cells(worksheet, cell_range))
        worksheet.merge_cells(str(cell_range))
    tables = spec.get("tables", [])
    if not isinstance(tables, list):
        fail("INVALID_ARGUMENT", "worksheet.tables 必须是数组")
    for table in tables:
        spreadsheet_add_table(worksheet, table)
    charts = spec.get("charts", [])
    if not isinstance(charts, list):
        fail("INVALID_ARGUMENT", "worksheet.charts 必须是数组")
    for chart in charts:
        spreadsheet_add_chart(worksheet, chart)
    if spec.get("hidden") is not None:
        worksheet.sheet_state = "hidden" if spec["hidden"] else "visible"


def workbook_from_spec(spec: dict[str, Any], template: Path | None = None) -> Workbook:
    if template is not None:
        workbook = load_workbook(str(template), data_only=False, keep_links=False)
    else:
        workbook = Workbook()
    worksheets = spec.get("worksheets", [])
    if not isinstance(worksheets, list):
        fail("INVALID_ARGUMENT", "spec.worksheets 必须是数组")
    if template is None and worksheets:
        default = workbook.active
        first = worksheets[0]
        if not isinstance(first, dict):
            fail("INVALID_ARGUMENT", "spec.worksheets 条目必须是对象")
        default.title = spreadsheet_sheet_name(first.get("name") or "Sheet1")
        apply_worksheet_spec(default, first)
        remaining = worksheets[1:]
    else:
        remaining = worksheets
    for worksheet_spec in remaining:
        if not isinstance(worksheet_spec, dict):
            fail("INVALID_ARGUMENT", "spec.worksheets 条目必须是对象")
        name = spreadsheet_sheet_name(worksheet_spec.get("name"))
        worksheet = workbook[name] if name in workbook.sheetnames else workbook.create_sheet(name)
        apply_worksheet_spec(worksheet, worksheet_spec)
    if len(workbook.worksheets) > MAX_WORKSHEETS:
        fail("RESOURCE_LIMIT", "XLSX 工作表数量超过限制")
    if not workbook.worksheets:
        workbook.create_sheet("Sheet1")
    configure_spreadsheet_calculation(workbook)
    return workbook


def xlsx_inspect(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".xlsx"})
    security = ooxml_preflight(source)
    workbook = load_workbook(str(source), data_only=False, keep_links=False)
    if len(workbook.worksheets) > MAX_WORKSHEETS:
        fail("RESOURCE_LIMIT", "XLSX 工作表数量超过限制")
    worksheets = []
    total_cells = 0
    for worksheet in workbook.worksheets:
        total_cells += worksheet.max_row * worksheet.max_column
        if total_cells > MAX_WORKSHEET_CELLS:
            fail("RESOURCE_LIMIT", "XLSX 已使用单元格数量超过限制")
        formulas = [
            {"cell": cell.coordinate, "formula": str(cell.value)[:2_000]}
            for row in worksheet.iter_rows()
            for cell in row
            if isinstance(cell.value, str) and cell.value.startswith("=")
        ]
        for formula in formulas:
            validate_spreadsheet_formula(formula["formula"])
        worksheets.append({
            "id": f"worksheet:{worksheet.title}",
            "name": worksheet.title,
            "state": worksheet.sheet_state,
            "rows": worksheet.max_row,
            "columns": worksheet.max_column,
            "freezePane": str(worksheet.freeze_panes or ""),
            "printArea": str(worksheet.print_area or ""),
            "autoFilter": str(worksheet.auto_filter.ref or ""),
            "mergedCells": [str(value) for value in list(worksheet.merged_cells.ranges)[:1_000]],
            "tables": [
                {
                    "name": table.displayName,
                    "ref": table.ref,
                    "style": table.tableStyleInfo.name if table.tableStyleInfo else None,
                }
                for table in worksheet.tables.values()
            ],
            "charts": [
                {
                    "id": f"chart:{worksheet.title}:{index}",
                    "index": index,
                    "type": chart.__class__.__name__,
                    "seriesCount": len(chart.series),
                }
                for index, chart in enumerate(worksheet._charts)
            ],
            "formulaCount": len(formulas),
            "formulas": formulas[:2_000],
            "preview": [
                [spreadsheet_value(cell.value) for cell in row[:20]]
                for row in list(
                    worksheet.iter_rows(
                        min_row=1,
                        max_row=min(worksheet.max_row, 20),
                        min_col=1,
                        max_col=min(worksheet.max_column, 20),
                    )
                )
            ],
        })
    return base_result(
        format="xlsx",
        sourcePath=relative_path(source),
        sourceHash=sha256_file(source),
        worksheetCount=len(worksheets),
        worksheets=worksheets,
        namedRanges=[
            {
                "name": value.name,
                "value": str(value.attr_text or ""),
            }
            for value in workbook.defined_names.values()
        ],
        security=security,
        warnings=[],
    )


def xlsx_create(payload: dict[str, Any]) -> dict[str, Any]:
    output = resolve_output(payload.get("outputPath"), ".xlsx")
    spec = payload.get("spec")
    if not isinstance(spec, dict):
        fail("INVALID_ARGUMENT", "spec 必须是对象")
    template_path = payload.get("templatePath")
    template = None
    template_id = None
    template_hash = None
    if template_path:
        template = resolve_source(template_path, "templatePath", {".xlsx"})
        ooxml_preflight(template)
        template_hash = sha256_file(template)
        if payload.get("templateHash") and template_hash != expected_hash(payload["templateHash"]):
            fail("SOURCE_CHANGED", "模板内容与 templateHash 不一致")
        template_id = Path(str(template_path)).stem
    workbook = workbook_from_spec(spec, template)
    atomic_save(workbook.save, output)
    recalculate_spreadsheet(output)
    ooxml_preflight(output)
    return base_result(
        artifact=artifact_record(
            output,
            metadata={
                "templateId": template_id,
                "templateHash": template_hash,
                "worksheetCount": len(workbook.worksheets),
                "recalculated": True,
            },
        ),
        warnings=[],
    )


def xlsx_edit(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".xlsx"})
    source_hash = verify_source_hash(source, payload.get("sourceHash"))
    ooxml_preflight(source)
    output = resolve_output(payload.get("outputPath"), ".xlsx")
    operations = payload.get("operations")
    if not isinstance(operations, list) or not operations:
        fail("INVALID_ARGUMENT", "operations 必须是非空数组")
    workbook = load_workbook(str(source), data_only=False, keep_links=False)
    results: list[dict[str, Any]] = []
    for index, operation in enumerate(operations):
        if not isinstance(operation, dict):
            fail("INVALID_ARGUMENT", f"operations[{index}] 必须是对象")
        name = str(operation.get("op") or "")
        if name == "set_cells":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            spreadsheet_set_cells(worksheet, operation.get("cells"))
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "set_range":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            spreadsheet_write_range(
                worksheet,
                operation.get("startCell") or "A1",
                operation.get("values"),
            )
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "add_worksheet":
            spec = operation.get("worksheet")
            if not isinstance(spec, dict):
                fail("INVALID_ARGUMENT", "add_worksheet.worksheet 必须是对象")
            sheet_name = spreadsheet_sheet_name(spec.get("name"))
            if sheet_name in workbook.sheetnames:
                fail("OUTPUT_EXISTS", f"工作表已存在：{sheet_name}")
            worksheet = workbook.create_sheet(sheet_name)
            apply_worksheet_spec(worksheet, spec)
            results.append({"index": index, "op": name, "sheet": sheet_name})
        elif name == "rename_worksheet":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            new_name = spreadsheet_sheet_name(operation.get("newName"))
            if new_name in workbook.sheetnames:
                fail("OUTPUT_EXISTS", f"工作表已存在：{new_name}")
            worksheet.title = new_name
            results.append({"index": index, "op": name, "sheet": new_name})
        elif name == "delete_worksheet":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            if len(workbook.worksheets) == 1:
                fail("VALIDATION_FAILED", "XLSX 至少需要一个工作表")
            workbook.remove(worksheet)
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "set_style":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            for cell in spreadsheet_range_cells(worksheet, operation.get("range")):
                spreadsheet_style(cell, operation.get("style"))
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "add_table":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            spreadsheet_add_table(worksheet, operation.get("table"))
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name in {"set_chart", "set_chart_data"}:
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            chart_value = operation.get("chart")
            if not isinstance(chart_value, dict):
                fail("INVALID_ARGUMENT", f"{name}.chart 必须是对象")
            chart_index = operation.get("chartIndex")
            if chart_index is not None:
                chart_index = int(chart_index)
                if chart_index < 0 or chart_index >= len(worksheet._charts):
                    fail("INVALID_ARGUMENT", "chartIndex 超出范围")
                worksheet._charts.pop(chart_index)
            spreadsheet_add_chart(worksheet, chart_value)
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "set_print_area":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            print_area = str(operation.get("range") or "")
            list(spreadsheet_range_cells(worksheet, print_area))
            worksheet.print_area = print_area
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "freeze_panes":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            worksheet.freeze_panes = str(operation.get("cell")) if operation.get("cell") else None
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        elif name == "set_column_widths":
            worksheet = spreadsheet_sheet(workbook, operation.get("sheet"))
            widths = operation.get("widths")
            if not isinstance(widths, dict):
                fail("INVALID_ARGUMENT", "set_column_widths.widths 必须是对象")
            for column, width in widths.items():
                key = get_column_letter(int(column)) if str(column).isdigit() else str(column).upper()
                worksheet.column_dimensions[key].width = float(width)
            results.append({"index": index, "op": name, "sheet": worksheet.title})
        else:
            fail("UNSUPPORTED_FEATURE", f"XLSX 不支持操作 {name}")
    if len(workbook.worksheets) > MAX_WORKSHEETS:
        fail("RESOURCE_LIMIT", "XLSX 工作表数量超过限制")
    configure_spreadsheet_calculation(workbook)
    atomic_save(workbook.save, output)
    recalculate_spreadsheet(output)
    ooxml_preflight(output)
    return base_result(
        artifact=artifact_record(
            output,
            metadata={
                "sourcePath": relative_path(source),
                "worksheetCount": len(workbook.worksheets),
                "recalculated": True,
            },
            source_hash=source_hash,
        ),
        operations=results,
        warnings=[],
    )


def spreadsheet_formula_errors(source: Path) -> dict[str, Any]:
    formulas = load_workbook(str(source), data_only=False, keep_links=False)
    values = load_workbook(str(source), data_only=True, keep_links=False)
    error_values = {"#REF!", "#VALUE!", "#DIV/0!", "#NAME?", "#NUM!", "#N/A", "#NULL!"}
    errors: list[dict[str, Any]] = []
    unresolved: list[dict[str, Any]] = []
    formula_count = 0
    for formula_sheet in formulas.worksheets:
        value_sheet = values[formula_sheet.title]
        for row in formula_sheet.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str) or not cell.value.startswith("="):
                    continue
                validate_spreadsheet_formula(cell.value)
                formula_count += 1
                cached = value_sheet[cell.coordinate].value
                if cached in error_values:
                    errors.append({
                        "sheet": formula_sheet.title,
                        "cell": cell.coordinate,
                        "formula": cell.value[:2_000],
                        "value": cached,
                    })
                elif cached is None:
                    unresolved.append({
                        "sheet": formula_sheet.title,
                        "cell": cell.coordinate,
                        "formula": cell.value[:2_000],
                    })
    return {
        "formulaCount": formula_count,
        "errors": errors[:2_000],
        "unresolved": unresolved[:2_000],
    }


def pdf_security(reader: PdfReader) -> dict[str, Any]:
    if reader.is_encrypted:
        fail("SECURITY_REJECTED", "PDF 已加密，V1 不处理加密文件")
    if len(reader.pages) > MAX_PDF_PAGES:
        fail("RESOURCE_LIMIT", "PDF 页数超过限制")
    root = reader.trailer.get("/Root", {})
    names = root.get("/Names") if hasattr(root, "get") else None
    javascript = bool(
        (hasattr(root, "get") and (root.get("/OpenAction") or root.get("/AA")))
        or (hasattr(names, "get") and names.get("/JavaScript"))
    )
    attachments = bool(hasattr(names, "get") and names.get("/EmbeddedFiles"))
    launch_actions = any(
        str(page.get("/AA") or "").find("/Launch") >= 0
        for page in reader.pages
    )
    if javascript or attachments or launch_actions:
        fail("SECURITY_REJECTED", "PDF 包含 JavaScript、Launch Action 或附件")
    return {
        "encrypted": False,
        "javascript": False,
        "attachments": False,
        "launchActions": False,
    }


def pdf_inspect(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions={".pdf"})
    reader = PdfReader(str(source), strict=True)
    security = pdf_security(reader)
    pages = []
    for index, page in enumerate(reader.pages):
        box = page.mediabox
        text = page.extract_text() or ""
        pages.append({
            "id": f"page:{index + 1}",
            "number": index + 1,
            "widthPoints": float(box.width),
            "heightPoints": float(box.height),
            "rotation": int(page.get("/Rotate") or 0),
            "textPreview": text[:4_000],
        })
    fields = reader.get_fields() or {}
    return base_result(
        format="pdf",
        sourcePath=relative_path(source),
        sourceHash=sha256_file(source),
        pageCount=len(reader.pages),
        pages=pages,
        forms=[
            {
                "name": name,
                "type": str(field.get("/FT") or ""),
                "value": field.get("/V"),
            }
            for name, field in list(fields.items())[:1_000]
        ],
        metadata={str(key): str(value)[:2_000] for key, value in (reader.metadata or {}).items()},
        security=security,
        warnings=[],
    )


def pdf_alignment(value: Any) -> int:
    return {"left": TA_LEFT, "center": TA_CENTER, "right": TA_RIGHT}.get(str(value), TA_LEFT)


def register_pdf_fonts() -> str:
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        return "STSong-Light"
    except Exception:
        return "Helvetica"


def pdf_paragraph(text: Any, style: ParagraphStyle) -> Paragraph:
    return Paragraph(escape(str(text or "")).replace("\n", "<br/>"), style)


def pdf_create(payload: dict[str, Any]) -> dict[str, Any]:
    output = resolve_output(payload.get("outputPath"), ".pdf")
    spec = payload.get("spec")
    if not isinstance(spec, dict):
        fail("INVALID_ARGUMENT", "spec 必须是对象")
    font_name = register_pdf_fonts()
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "MeteoMateBody",
        parent=styles["BodyText"],
        fontName=font_name,
        fontSize=float(spec.get("fontSize") or 10.5),
        leading=float(spec.get("leading") or 17),
        wordWrap="CJK",
        alignment=pdf_alignment(spec.get("alignment")),
        spaceAfter=6,
    )
    heading = ParagraphStyle(
        "MeteoMateHeading",
        parent=styles["Heading1"],
        fontName=font_name,
        fontSize=18,
        leading=25,
        wordWrap="CJK",
        alignment=TA_CENTER,
        spaceAfter=14,
    )
    page_size = LETTER if str(spec.get("pageSize") or "A4").upper() == "LETTER" else A4
    if str(spec.get("orientation") or "").lower() == "landscape":
        page_size = landscape(page_size)
    margin = float(spec.get("marginPoints") or 54)
    story: list[Any] = []
    if spec.get("title"):
        story.append(pdf_paragraph(spec["title"], heading))
    blocks = spec.get("blocks", [])
    if not isinstance(blocks, list):
        fail("INVALID_ARGUMENT", "spec.blocks 必须是数组")
    for block in blocks:
        if not isinstance(block, dict):
            fail("INVALID_ARGUMENT", "PDF 内容块必须是对象")
        kind = str(block.get("type") or "paragraph")
        if kind == "paragraph":
            story.append(pdf_paragraph(block.get("text"), body))
        elif kind == "heading":
            level = max(1, min(3, int(block.get("level") or 1)))
            local = ParagraphStyle(
                f"MeteoMateHeading{level}",
                parent=body,
                fontSize={1: 16, 2: 14, 3: 12}[level],
                leading={1: 22, 2: 20, 3: 18}[level],
                spaceBefore=8,
                spaceAfter=6,
            )
            story.append(pdf_paragraph(block.get("text"), local))
        elif kind == "table":
            rows = block.get("rows")
            if not isinstance(rows, list) or not rows or not all(isinstance(row, list) for row in rows):
                fail("INVALID_ARGUMENT", "table.rows 必须是非空二维数组")
            cells = [[pdf_paragraph(value, body) for value in row] for row in rows]
            table = Table(cells, repeatRows=1 if block.get("header", True) else 0)
            table.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, -1), font_name),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#94A3B8")),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E2E8F0")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 5),
                ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]))
            story.append(table)
            story.append(Spacer(1, 8))
        elif kind == "image":
            image_path = resolve_source(block.get("path"), "image.path")
            validate_image(image_path)
            width = float(block.get("widthPoints") or 420)
            with Image.open(image_path) as image:
                ratio = image.height / image.width
            story.append(ReportLabImage(str(image_path), width=width, height=width * ratio))
            story.append(Spacer(1, 8))
        elif kind == "page_break":
            story.append(PageBreak())
        elif kind == "spacer":
            story.append(Spacer(1, float(block.get("heightPoints") or 12)))
        else:
            fail("UNSUPPORTED_FEATURE", f"PDF 不支持内容块 {kind}")
    if not story:
        story.append(pdf_paragraph("", body))

    header = str(spec.get("header") or "")
    footer = str(spec.get("footer") or "")

    def decorate_page(canvas: Any, document: Any) -> None:
        canvas.saveState()
        canvas.setFont(font_name, 9)
        if header:
            canvas.drawString(document.leftMargin, page_size[1] - 28, header)
        footer_text = footer.replace("{page}", str(document.page))
        if footer_text:
            canvas.drawString(document.leftMargin, 24, footer_text)
        canvas.restoreState()

    def save_pdf(target: Path) -> None:
        document = SimpleDocTemplate(
            str(target),
            pagesize=page_size,
            leftMargin=margin,
            rightMargin=margin,
            topMargin=margin,
            bottomMargin=margin,
            title=str(spec.get("title") or ""),
            author="MeteoMate",
        )
        document.build(story, onFirstPage=decorate_page, onLaterPages=decorate_page)

    atomic_save(save_pdf, output)
    pdf_security(PdfReader(str(output), strict=True))
    return base_result(artifact=artifact_record(output), warnings=[])


def watermark_page(text: str, width: float, height: float, operation: dict[str, Any]) -> Any:
    from reportlab.pdfgen import canvas

    stream = io.BytesIO()
    layer = canvas.Canvas(stream, pagesize=(width, height))
    font_name = register_pdf_fonts()
    layer.saveState()
    if hasattr(layer, "setFillAlpha"):
        layer.setFillAlpha(float(operation.get("opacity") or 0.18))
    layer.setFillColor(colors.HexColor(str(operation.get("color") or "#64748B")))
    layer.setFont(font_name, float(operation.get("fontSize") or 42))
    x = float(operation.get("x") or width / 2)
    y = float(operation.get("y") or height / 2)
    layer.translate(x, y)
    layer.rotate(float(operation.get("rotation") or 35))
    layer.drawCentredString(0, 0, text)
    layer.restoreState()
    layer.save()
    stream.seek(0)
    return PdfReader(stream).pages[0]


def sanitize_pdf_writer(writer: PdfWriter) -> None:
    root = writer.root_object
    root.pop("/OpenAction", None)
    root.pop("/AA", None)
    names = root.get("/Names")
    if hasattr(names, "pop"):
        names.pop("/JavaScript", None)
        names.pop("/EmbeddedFiles", None)


def pdf_transform(payload: dict[str, Any]) -> dict[str, Any]:
    input_values = payload.get("inputs")
    if not isinstance(input_values, list) or not input_values:
        fail("INVALID_ARGUMENT", "inputs 必须是非空数组")
    sources = [resolve_source(value, f"inputs[{index}]", {".pdf"}) for index, value in enumerate(input_values)]
    output = resolve_output(payload.get("outputPath"), ".pdf")
    readers = [PdfReader(str(source), strict=True) for source in sources]
    for reader in readers:
        pdf_security(reader)
    writer = PdfWriter()
    for reader in readers:
        writer.append(reader)
    sanitize_pdf_writer(writer)
    operations = payload.get("operations")
    if not isinstance(operations, list) or not operations:
        fail("INVALID_ARGUMENT", "operations 必须是非空数组")
    results: list[dict[str, Any]] = []
    for index, operation in enumerate(operations):
        if not isinstance(operation, dict):
            fail("INVALID_ARGUMENT", f"operations[{index}] 必须是对象")
        name = str(operation.get("op") or "")
        if name == "merge":
            results.append({"index": index, "op": name, "pageCount": len(writer.pages)})
        elif name == "split":
            pages = operation.get("pages")
            if not isinstance(pages, list) or not pages:
                fail("INVALID_ARGUMENT", "split.pages 必须是页码数组")
            selected = PdfWriter()
            for page_number in pages:
                number = int(page_number)
                if number < 1 or number > len(writer.pages):
                    fail("INVALID_ARGUMENT", f"页码超出范围：{number}")
                selected.add_page(copy.copy(writer.pages[number - 1]))
            writer = selected
            results.append({"index": index, "op": name, "pageCount": len(writer.pages)})
        elif name == "remove_pages":
            removed = {int(value) for value in operation.get("pages", [])}
            selected = PdfWriter()
            for page_number, page in enumerate(writer.pages, start=1):
                if page_number not in removed:
                    selected.add_page(copy.copy(page))
            writer = selected
            results.append({"index": index, "op": name, "pageCount": len(writer.pages)})
        elif name == "rotate":
            degrees = int(operation.get("degrees") or 0)
            if degrees % 90:
                fail("INVALID_ARGUMENT", "rotate.degrees 必须是 90 的倍数")
            selected_pages = {int(value) for value in operation.get("pages", [])}
            for page_number, page in enumerate(writer.pages, start=1):
                if not selected_pages or page_number in selected_pages:
                    page.rotate(degrees)
            results.append({"index": index, "op": name, "degrees": degrees})
        elif name == "watermark":
            text = str(operation.get("text") or "")
            if not text:
                fail("INVALID_ARGUMENT", "watermark.text 不能为空")
            selected_pages = {int(value) for value in operation.get("pages", [])}
            for page_number, page in enumerate(writer.pages, start=1):
                if selected_pages and page_number not in selected_pages:
                    continue
                layer = watermark_page(text, float(page.mediabox.width), float(page.mediabox.height), operation)
                page.merge_page(layer)
            results.append({"index": index, "op": name})
        elif name == "fill_form":
            fields = operation.get("fields")
            if not isinstance(fields, dict) or not fields:
                fail("INVALID_ARGUMENT", "fill_form.fields 必须是非空对象")
            for page in writer.pages:
                try:
                    writer.update_page_form_field_values(page, fields, auto_regenerate=False)
                except Exception as error:
                    fail("VALIDATION_FAILED", f"PDF 表单填写失败：{error}")
            results.append({"index": index, "op": name, "fields": sorted(fields)})
        elif name == "add_blank_page":
            writer.add_blank_page(
                width=float(operation.get("widthPoints") or A4[0]),
                height=float(operation.get("heightPoints") or A4[1]),
            )
            results.append({"index": index, "op": name})
        else:
            fail("UNSUPPORTED_FEATURE", f"PDF 不支持操作 {name}")
    if not writer.pages:
        fail("VALIDATION_FAILED", "PDF 转换结果没有页面")
    sanitize_pdf_writer(writer)

    def save_pdf(target: Path) -> None:
        with target.open("wb") as handle:
            writer.write(handle)

    atomic_save(save_pdf, output)
    pdf_security(PdfReader(str(output), strict=True))
    return base_result(
        artifact=artifact_record(
            output,
            metadata={"sourcePaths": [relative_path(source) for source in sources]},
        ),
        operations=results,
        warnings=[],
    )


def soffice_path() -> Path:
    raw = os.environ.get("METEOMATE_SOFFICE_PATH", "")
    if not raw:
        fail("RUNTIME_UNAVAILABLE", "Office 渲染与公式重算需要 Managed Runtime 中的 LibreOffice")
    target = Path(raw)
    if not target.is_absolute() or not target.is_file() or not os.access(target, os.X_OK):
        fail("RUNTIME_UNAVAILABLE", "LibreOffice 入口不可用")
    return target


def libreoffice_convert(
    source: Path,
    output_directory: Path,
    output_format: str,
    timeout: int = 120,
) -> Path:
    executable = soffice_path()
    profile = Path(tempfile.mkdtemp(prefix="meteomate-lo-profile-"))
    try:
        profile_uri = profile.resolve().as_uri()
        result = subprocess.run(
            [
                str(executable),
                f"-env:UserInstallation={profile_uri}",
                "--headless",
                "--nologo",
                "--nodefault",
                "--nofirststartwizard",
                "--norestore",
                "--convert-to",
                output_format,
                "--outdir",
                str(output_directory),
                str(source),
            ],
            cwd=output_directory,
            env={
                "PATH": os.environ.get("PATH", ""),
                "LANG": os.environ.get("LANG", "C.UTF-8"),
                "HOME": str(profile),
            },
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
        )
        extension = output_format.split(":", 1)[0].lower()
        converted = output_directory / f"{source.stem}.{extension}"
        if result.returncode != 0 or not converted.exists():
            message = (result.stderr or result.stdout or "").strip().splitlines()
            fail(
                "CONVERSION_FAILED",
                message[-1] if message else f"LibreOffice 未生成 {extension.upper()}",
            )
        os.chmod(converted, stat.S_IRUSR | stat.S_IWUSR)
        return converted
    except subprocess.TimeoutExpired:
        fail("CONVERSION_FAILED", f"LibreOffice 转换超过 {timeout} 秒")
    finally:
        shutil.rmtree(profile, ignore_errors=True)


def recalculate_spreadsheet(source: Path) -> None:
    temporary_root = Path(tempfile.mkdtemp(prefix="meteomate-xlsx-recalculate-"))
    input_directory = temporary_root / "input"
    output_directory = temporary_root / "output"
    input_directory.mkdir(mode=0o700)
    output_directory.mkdir(mode=0o700)
    temporary_input = input_directory / source.name
    shutil.copyfile(source, temporary_input)
    try:
        recalculated = libreoffice_convert(
            temporary_input,
            output_directory,
            "xlsx",
            timeout=120,
        )
        ooxml_preflight(recalculated)
        temporary_output = source.with_name(f".{source.name}.recalculated-{os.getpid()}")
        shutil.copyfile(recalculated, temporary_output)
        os.chmod(temporary_output, stat.S_IRUSR | stat.S_IWUSR)
        temporary_output.replace(source)
    finally:
        shutil.rmtree(temporary_root, ignore_errors=True)


def preview_directory(source_hash: str) -> Path:
    private_root = WORKSPACE / ".meteomate"
    private_root.mkdir(parents=True, exist_ok=True, mode=0o700)
    os.chmod(private_root, stat.S_IRWXU)
    previews_root = private_root / "previews"
    previews_root.mkdir(exist_ok=True, mode=0o700)
    os.chmod(previews_root, stat.S_IRWXU)
    target = previews_root / source_hash[:24]
    target.mkdir(parents=True, exist_ok=True, mode=0o700)
    os.chmod(target, stat.S_IRWXU)
    if not inside_workspace(target.resolve(strict=True)):
        fail("WORKSPACE_VIOLATION", "预览目录已超出项目工作区")
    return target


def write_json(target: Path, value: Any) -> None:
    temporary = target.with_name(f".{target.name}.tmp-{os.getpid()}")
    temporary.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    os.chmod(temporary, stat.S_IRUSR | stat.S_IWUSR)
    temporary.replace(target)


def render_internal(
    source: Path,
    *,
    dpi: int = 144,
    page_from: int = 1,
    page_to: int | None = None,
) -> dict[str, Any]:
    source_hash = sha256_file(source)
    output_directory = preview_directory(source_hash)
    temporary_directory = Path(tempfile.mkdtemp(prefix="meteomate-office-render-"))
    try:
        if source.suffix.lower() in {".docx", ".pptx", ".xlsx"}:
            ooxml_preflight(source)
            converted = libreoffice_convert(source, temporary_directory, "pdf")
        elif source.suffix.lower() == ".pdf":
            reader = PdfReader(str(source), strict=True)
            pdf_security(reader)
            converted = source
        else:
            fail("INVALID_ARGUMENT", "artifact_render 仅支持 DOCX、PPTX、XLSX 和 PDF")
        preview_pdf = output_directory / "preview.pdf"
        if converted.resolve() != preview_pdf.resolve():
            temporary_pdf = output_directory / f".preview.pdf.tmp-{os.getpid()}"
            shutil.copyfile(converted, temporary_pdf)
            os.chmod(temporary_pdf, stat.S_IRUSR | stat.S_IWUSR)
            temporary_pdf.replace(preview_pdf)

        import pypdfium2 as pdfium

        document = pdfium.PdfDocument(str(preview_pdf))
        page_count = len(document)
        if page_count > MAX_RENDER_PAGES:
            fail("RESOURCE_LIMIT", "预览页数超过限制")
        last_page = min(page_to or page_count, page_count)
        first_page = max(1, page_from)
        if first_page > last_page:
            fail("INVALID_ARGUMENT", "pages.from 不能大于 pages.to")
        thumbnails: list[str] = []
        scale = dpi / 72
        for page_number in range(first_page, last_page + 1):
            page = document[page_number - 1]
            bitmap = page.render(scale=scale)
            image = bitmap.to_pil()
            target = output_directory / f"page-{page_number}.png"
            image.save(target, format="PNG", optimize=True)
            os.chmod(target, stat.S_IRUSR | stat.S_IWUSR)
            thumbnails.append(relative_path(target))
            image.close()
            bitmap.close()
            page.close()
        document.close()
        manifest = {
            "schemaVersion": "meteomate.preview/v1",
            "sourcePath": relative_path(source),
            "sourceHash": source_hash,
            "previewPath": relative_path(preview_pdf),
            "pageCount": page_count,
            "dpi": dpi,
            "thumbnails": thumbnails,
        }
        manifest_path = output_directory / "manifest.json"
        write_json(manifest_path, manifest)
        return {
            **manifest,
            "previewManifestPath": relative_path(manifest_path),
            "thumbnailPath": thumbnails[0] if thumbnails else None,
        }
    finally:
        shutil.rmtree(temporary_directory, ignore_errors=True)


def artifact_render(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions=ALLOWED_EXTENSIONS)
    pages = payload.get("pages") if isinstance(payload.get("pages"), dict) else {}
    render = render_internal(
        source,
        dpi=max(72, min(240, int(payload.get("dpi") or 144))),
        page_from=max(1, int(pages.get("from") or 1)),
        page_to=int(pages["to"]) if pages.get("to") is not None else None,
    )
    return base_result(
        artifact=artifact_record(source, status="draft", metadata={"render": render}),
        render=render,
        warnings=[],
    )


def artifact_validate(payload: dict[str, Any]) -> dict[str, Any]:
    source = resolve_source(payload.get("sourcePath"), extensions=ALLOWED_EXTENSIONS)
    checks: list[dict[str, Any]] = []
    warnings: list[str] = []
    render: dict[str, Any] | None = None
    valid = True
    try:
        if source.suffix.lower() == ".docx":
            security = ooxml_preflight(source)
            Document(str(source))
            checks.append({"name": "structure", "status": "passed"})
            checks.append({"name": "security", "status": "passed", "detail": security})
        elif source.suffix.lower() == ".pptx":
            security = ooxml_preflight(source)
            presentation = Presentation(str(source))
            if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
                fail("RESOURCE_LIMIT", "PPTX 幻灯片数量超过限制")
            shape_count = sum(len(slide.shapes) for slide in presentation.slides)
            if shape_count > MAX_PRESENTATION_SHAPES:
                fail("RESOURCE_LIMIT", "PPTX 形状数量超过限制")
            checks.append({
                "name": "structure",
                "status": "passed",
                "slideCount": len(presentation.slides),
                "shapeCount": shape_count,
            })
            checks.append({"name": "security", "status": "passed", "detail": security})
        elif source.suffix.lower() == ".xlsx":
            security = ooxml_preflight(source)
            workbook = load_workbook(str(source), data_only=False, keep_links=False)
            if len(workbook.worksheets) > MAX_WORKSHEETS:
                fail("RESOURCE_LIMIT", "XLSX 工作表数量超过限制")
            formulas = spreadsheet_formula_errors(source)
            if formulas["errors"]:
                fail(
                    "VALIDATION_FAILED",
                    f"XLSX 包含 {len(formulas['errors'])} 个公式错误",
                )
            if formulas["unresolved"]:
                warnings.append(f"{len(formulas['unresolved'])} 个公式没有可读取的缓存结果")
            checks.append({
                "name": "structure",
                "status": "passed",
                "worksheetCount": len(workbook.worksheets),
            })
            checks.append({
                "name": "formulas",
                "status": "passed",
                **formulas,
            })
            checks.append({"name": "security", "status": "passed", "detail": security})
        else:
            reader = PdfReader(str(source), strict=True)
            security = pdf_security(reader)
            checks.append({"name": "structure", "status": "passed", "pageCount": len(reader.pages)})
            checks.append({"name": "security", "status": "passed", "detail": security})
        if payload.get("requireRender", True):
            render = render_internal(source)
            checks.append({"name": "render", "status": "passed", "pageCount": render["pageCount"]})
    except OfficeError as error:
        valid = False
        checks.append({"name": "validation", "status": "failed", "code": error.code, "message": str(error)})
    except Exception as error:
        valid = False
        checks.append({"name": "validation", "status": "failed", "code": "VALIDATION_FAILED", "message": str(error)})
    status = "ready" if valid else "failed"
    metadata: dict[str, Any] = {
        "validation": {
            "status": status,
            "checks": checks,
        },
    }
    if render:
        metadata["render"] = render
    return base_result(
        valid=valid,
        status=status,
        artifact=artifact_record(source, status=status, metadata=metadata),
        checks=checks,
        warnings=warnings,
        **({"previewManifestPath": render["previewManifestPath"]} if render else {}),
    )


TOOLS = {
    "docx_inspect": docx_inspect,
    "docx_create": docx_create,
    "docx_edit": docx_edit,
    "pptx_inspect": pptx_inspect,
    "pptx_create": pptx_create,
    "pptx_edit": pptx_edit,
    "xlsx_inspect": xlsx_inspect,
    "xlsx_create": xlsx_create,
    "xlsx_edit": xlsx_edit,
    "pdf_inspect": pdf_inspect,
    "pdf_create": pdf_create,
    "pdf_transform": pdf_transform,
    "artifact_render": artifact_render,
    "artifact_validate": artifact_validate,
}


def main() -> int:
    if len(sys.argv) != 2 or sys.argv[1] not in TOOLS:
        print("INVALID_ARGUMENT: 未知 Office 工具", file=sys.stderr)
        return 2
    try:
        payload = json.load(sys.stdin)
        if not isinstance(payload, dict):
            fail("INVALID_ARGUMENT", "工具输入必须是 JSON 对象")
        if payload.get("schemaVersion") != SCHEMA_VERSION:
            fail("INVALID_ARGUMENT", f"schemaVersion 必须是 {SCHEMA_VERSION}")
        result = TOOLS[sys.argv[1]](payload)
        sys.stdout.write(json.dumps(result, ensure_ascii=False, separators=(",", ":")))
        return 0
    except OfficeError as error:
        print(str(error), file=sys.stderr)
        return 1
    except Exception as error:
        print(f"OFFICE_TOOL_FAILED: {error}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
